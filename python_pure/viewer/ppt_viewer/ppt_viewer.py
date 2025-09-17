"""
Simple PPTX viewer using python-pptx + PyQt5 (no external dependencies).
Requirements:
  - Python packages: pyqt5, python-pptx
    pip install pyqt5 python-pptx

Usage:
  python pptx_viewer.py

Features:
  - Open a .pptx file (not .ppt)
  - Extract text and images from each slide
  - View slides with Prev / Next buttons
  - Independent of Microsoft PowerPoint / WPS / LibreOffice
Limitations:
  - Only supports PPTX (not old PPT)
  - No slide layout rendering (just sequential text + images)
"""

import sys
import os
import tempfile
from PyQt5.QtWidgets import (
    QApplication, QWidget, QLabel, QPushButton, QHBoxLayout, QVBoxLayout,
    QFileDialog, QMessageBox, QScrollArea, QSizePolicy
)
from PyQt5.QtGui import QPixmap
from PyQt5.QtCore import Qt
from pptx import Presentation

class PPTXViewer(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle('Simple PPTX Viewer (python-pptx)')
        self.resize(800, 600)

        self.slides_content = []
        self.current_index = -1
        self.temp_dir = tempfile.mkdtemp(prefix="pptxview_")

        # central widget with scroll
        self.scroll = QScrollArea()
        self.content_widget = QWidget()
        self.content_layout = QVBoxLayout()
        self.content_widget.setLayout(self.content_layout)
        self.scroll.setWidget(self.content_widget)
        self.scroll.setWidgetResizable(True)

        self.open_btn = QPushButton('Open')
        self.prev_btn = QPushButton('Prev')
        self.next_btn = QPushButton('Next')

        self.open_btn.clicked.connect(self.open_file)
        self.prev_btn.clicked.connect(self.prev_slide)
        self.next_btn.clicked.connect(self.next_slide)

        h = QHBoxLayout()
        h.addWidget(self.open_btn)
        h.addStretch(1)
        h.addWidget(self.prev_btn)
        h.addWidget(self.next_btn)

        v = QVBoxLayout()
        v.addWidget(self.scroll, 1)
        v.addLayout(h)
        self.setLayout(v)

    def open_file(self):
        path, _ = QFileDialog.getOpenFileName(self, 'Open Presentation', '', 'PowerPoint Files (*.pptx)')
        if not path:
            return
        try:
            prs = Presentation(path)
        except Exception as e:
            QMessageBox.critical(self, 'Error', f'Failed to open PPTX: {e}')
            return

        # clear old temp images
        for f in os.listdir(self.temp_dir):
            try:
                os.remove(os.path.join(self.temp_dir, f))
            except Exception:
                pass

        self.slides_content = []
        for idx, slide in enumerate(prs.slides, start=1):
            elements = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        elements.append(("text", shape.text.strip()))
                    if getattr(shape, "shape_type", None) == 13 and hasattr(shape, "image"):
                        image = shape.image
                        img_ext = image.ext
                        img_path = os.path.join(self.temp_dir, f"slide{idx}_{len(elements)}.{img_ext}")
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        elements.append(("image", img_path))
                except Exception:
                    continue
            self.slides_content.append(elements)

        if not self.slides_content:
            QMessageBox.warning(self, 'No content', 'No text or images found in this PPTX.')
            return

        self.current_index = 0
        self.show_current()

    def show_current(self):
        if self.current_index < 0 or self.current_index >= len(self.slides_content):
            return

        while self.content_layout.count():
            item = self.content_layout.takeAt(0)
            w = item.widget()
            if w:
                w.deleteLater()

        elements = self.slides_content[self.current_index]
        if not elements:
            lbl = QLabel('(empty slide)')
            lbl.setAlignment(Qt.AlignCenter)
            self.content_layout.addWidget(lbl)
        else:
            for kind, data in elements:
                if kind == "text":
                    lbl = QLabel(data)
                    lbl.setWordWrap(True)
                    lbl.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Minimum)
                    self.content_layout.addWidget(lbl)
                elif kind == "image":
                    pix = QPixmap(data)
                    img_lbl = QLabel()
                    img_lbl.setPixmap(pix.scaledToWidth(700, Qt.SmoothTransformation))
                    img_lbl.setAlignment(Qt.AlignCenter)
                    self.content_layout.addWidget(img_lbl)
        self.content_layout.addStretch(1)

    def prev_slide(self):
        if not self.slides_content:
            return
        if self.current_index > 0:
            self.current_index -= 1
            self.show_current()

    def next_slide(self):
        if not self.slides_content:
            return
        if self.current_index < len(self.slides_content) - 1:
            self.current_index += 1
            self.show_current()

    def closeEvent(self, event):
        try:
            for f in os.listdir(self.temp_dir):
                os.remove(os.path.join(self.temp_dir, f))
            os.rmdir(self.temp_dir)
        except Exception:
            pass
        super().closeEvent(event)

if __name__ == '__main__':
    app = QApplication(sys.argv)
    w = PPTXViewer()
    w.show()
    sys.exit(app.exec_())
